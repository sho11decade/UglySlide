#!/usr/bin/env python
"""Test: Save to BytesIO first, then to disk to ensure proper file writing"""

import os
import sys
from io import BytesIO
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from src.taco_generator import TacoGenerator
from src.content_transformer import ContentTransformer

print("=== Testing BytesIO → Disk write method ===\n")

test_file = 'samples/test_sample.pptx'

for level in [1, 5, 7, 10]:
    print(f"Level {level}:")
    
    try:
        # Load original
        prs = Presentation(test_file)
        
        # Apply transformations
        TacoGenerator(prs, tacky_level=level, seed=42).apply_tacky_design()
        ContentTransformer(prs, intensity=level, seed=42).transform_all_content()
        
        # Method 1: Direct save to file (original)
        output1 = f'uploads/direct_level_{level}.pptx'
        prs.save(output1)
        size1 = os.path.getsize(output1)
        print(f"  Direct save: {size1} bytes")
        
        # Verify direct
        Presentation(output1)
        print(f"  Direct ✓ Verified")
        
        # Method 2: Save to BytesIO first, then to file
        prs2 = Presentation(test_file)
        TacoGenerator(prs2, tacky_level=level, seed=42).apply_tacky_design()
        ContentTransformer(prs2, intensity=level, seed=42).transform_all_content()
        
        buffer = BytesIO()
        prs2.save(buffer)
        buffer_data = buffer.getvalue()
        
        output2 = f'uploads/buffered_level_{level}.pptx'
        with open(output2, 'wb') as f:
            f.write(buffer_data)
        
        size2 = os.path.getsize(output2)
        print(f"  Buffered save: {size2} bytes")
        
        # Verify buffered
        Presentation(output2)
        print(f"  Buffered ✓ Verified")
        
        print()
        
    except Exception as e:
        print(f"  ✗ FAILED: {e}\n")
        import traceback
        traceback.print_exc()
        break

print("--- Test Complete ---")
