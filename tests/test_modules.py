"""Test suite for DasaMaker"""

import pytest
import os
import sys
from pathlib import Path
from tempfile import TemporaryDirectory

# Add src to path
sys.path.insert(0, str(Path(__file__).parent.parent))

from src.ppt_analyzer import PPTAnalyzer, DesignAnalysis
from src.taco_generator import TacoGenerator
from src.content_transformer import ContentTransformer

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class TestDesignAnalysis:
    """Test DesignAnalysis class"""
    
    def test_design_analysis_initialization(self):
        """Test DesignAnalysis initialization"""
        analysis = DesignAnalysis()
        assert analysis.fonts == {}
        assert analysis.colors == []
        assert analysis.total_slides == 0
        assert analysis.animation_count == 0
    
    def test_design_analysis_to_dict(self):
        """Test converting analysis to dictionary"""
        analysis = DesignAnalysis()
        analysis.total_slides = 3
        analysis.fonts = {"Arial": 5, "Times": 3}
        analysis.colors = [(255, 0, 0), (0, 255, 0)]
        
        result = analysis.to_dict()
        assert result['total_slides'] == 3
        assert result['fonts'] == {"Arial": 5, "Times": 3}
        assert len(result['colors']) == 2


class TestPPTAnalyzer:
    """Test PPT Analyzer"""
    
    def create_test_pptx(self, temp_path: str) -> str:
        """Create a test PPTX file"""
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        
        # Add a slide
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add a text box
        left = top = width = height = Inches(1)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = "Test Slide"
        
        # Add a shape with fill
        shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = Rectangle
        shape.fill.solid()
        
        # Save PPTX
        output_path = os.path.join(temp_path, "test.pptx")
        prs.save(output_path)
        return output_path
    
    def test_ppt_analyzer_initialization(self):
        """Test PPTAnalyzer initialization with valid file"""
        with TemporaryDirectory() as temp_dir:
            test_file = self.create_test_pptx(temp_dir)
            analyzer = PPTAnalyzer(test_file)
            assert analyzer.pptx_path.exists()
            assert analyzer.presentation is not None
    
    def test_ppt_analyzer_invalid_file(self):
        """Test PPTAnalyzer with invalid file"""
        with pytest.raises(FileNotFoundError):
            PPTAnalyzer("nonexistent.pptx")
    
    def test_ppt_analyzer_invalid_extension(self):
        """Test PPTAnalyzer with invalid file extension"""
        with TemporaryDirectory() as temp_dir:
            txt_file = os.path.join(temp_dir, "test.txt")
            Path(txt_file).write_text("test")
            
            with pytest.raises(ValueError):
                PPTAnalyzer(txt_file)
    
    def test_ppt_analyzer_analyze(self):
        """Test analyzing a presentation"""
        with TemporaryDirectory() as temp_dir:
            test_file = self.create_test_pptx(temp_dir)
            analyzer = PPTAnalyzer(test_file)
            analysis = analyzer.analyze()
            
            assert analysis.total_slides == 1
            assert isinstance(analysis.fonts, dict)
            assert isinstance(analysis.colors, list)
    
    def test_ppt_analyzer_dominant_fonts(self):
        """Test getting dominant fonts"""
        with TemporaryDirectory() as temp_dir:
            test_file = self.create_test_pptx(temp_dir)
            analyzer = PPTAnalyzer(test_file)
            fonts = analyzer.get_dominant_fonts(top_n=5)
            
            assert isinstance(fonts, list)
            # Check all items are (name, count) tuples
            for font_name, count in fonts:
                assert isinstance(font_name, str)
                assert isinstance(count, int)


class TestTacoGenerator:
    """Test Taco Generator"""
    
    def create_test_presentation(self) -> Presentation:
        """Create a test presentation"""
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add text
        left = top = width = height = Inches(1)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Professional Text"
        
        return prs
    
    def test_taco_generator_initialization(self):
        """Test TacoGenerator initialization"""
        prs = self.create_test_presentation()
        generator = TacoGenerator(prs, tacky_level=5)
        
        assert generator.tacky_level == 5
        assert generator.presentation == prs
    
    def test_taco_generator_level_clamping(self):
        """Test that tackiness level is clamped to 1-10"""
        prs = self.create_test_presentation()
        
        gen_low = TacoGenerator(prs, tacky_level=-5)
        assert gen_low.tacky_level == 1
        
        gen_high = TacoGenerator(prs, tacky_level=15)
        assert gen_high.tacky_level == 10
    
    def test_taco_generator_apply_design(self):
        """Test applying tacky design"""
        prs = self.create_test_presentation()
        generator = TacoGenerator(prs, tacky_level=7)
        
        # Should not raise an error
        generator.apply_tacky_design()
    
    def test_taco_generator_set_level(self):
        """Test setting tackiness level"""
        prs = self.create_test_presentation()
        generator = TacoGenerator(prs, tacky_level=5)
        
        generator.set_tacky_level(8)
        assert generator.get_tacky_level() == 8
    
    def test_taco_generator_gradient_level_5(self):
        """Test gradient application at level 5"""
        prs = self.create_test_presentation()
        generator = TacoGenerator(prs, tacky_level=5, seed=42)
        
        # Should not raise an error
        generator.apply_tacky_design()
    
    def test_taco_generator_gradient_level_7(self):
        """Test gradient application at level 7"""
        prs = self.create_test_presentation()
        generator = TacoGenerator(prs, tacky_level=7, seed=42)
        
        # Should not raise an error
        generator.apply_tacky_design()
    
    def test_taco_generator_gradient_level_10(self):
        """Test extreme gradient application at level 10"""
        prs = self.create_test_presentation()
        generator = TacoGenerator(prs, tacky_level=10, seed=42)
        
        # Should not raise an error
        generator.apply_tacky_design()
    
    def test_taco_generator_with_shapes(self):
        """Test tacky generator with shapes that have fills"""
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add a shape with fill
        left = top = Inches(1)
        width = height = Inches(2)
        shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 0, 0)
        
        generator = TacoGenerator(prs, tacky_level=7, seed=42)
        generator.apply_tacky_design()


class TestContentTransformer:
    """Test Content Transformer"""
    
    def create_test_presentation(self) -> Presentation:
        """Create a test presentation"""
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add text
        left = top = width = height = Inches(1)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "This is a test bullet point."
        
        return prs
    
    def test_content_transformer_initialization(self):
        """Test ContentTransformer initialization"""
        prs = self.create_test_presentation()
        transformer = ContentTransformer(prs, intensity=5)
        
        assert transformer.intensity == 5
        assert transformer.presentation == prs
    
    def test_content_transformer_intensity_clamping(self):
        """Test that intensity is clamped to 1-10"""
        prs = self.create_test_presentation()
        
        trans_low = ContentTransformer(prs, intensity=-5)
        assert trans_low.intensity == 1
        
        trans_high = ContentTransformer(prs, intensity=15)
        assert trans_high.intensity == 10
    
    def test_content_transformer_transform(self):
        """Test applying content transformation"""
        prs = self.create_test_presentation()
        transformer = ContentTransformer(prs, intensity=7)
        
        # Should not raise an error
        transformer.transform_all_content()
    
    def test_content_transformer_set_intensity(self):
        """Test setting transformation intensity"""
        prs = self.create_test_presentation()
        transformer = ContentTransformer(prs, intensity=5)
        
        transformer.set_intensity(8)
        assert transformer.get_intensity() == 8


class TestIntegration:
    """Integration tests"""
    
    def create_test_pptx(self, temp_path: str) -> str:
        """Create a test PPTX file"""
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        
        # Add multiple slides
        for i in range(3):
            slide = prs.slides.add_slide(blank_slide_layout)
            left = top = width = height = Inches(1)
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = f"Slide {i + 1} Content"
        
        output_path = os.path.join(temp_path, "test.pptx")
        prs.save(output_path)
        return output_path
    
    def test_full_pipeline(self):
        """Test the full processing pipeline"""
        with TemporaryDirectory() as temp_dir:
            # Create test file
            input_file = self.create_test_pptx(temp_dir)
            
            # Step 1: Analyze
            analyzer = PPTAnalyzer(input_file)
            analysis = analyzer.analyze()
            assert analysis.total_slides == 3
            
            # Step 2: Load for modification
            presentation = Presentation(input_file)
            
            # Step 3: Apply tacky design
            generator = TacoGenerator(presentation, tacky_level=7)
            generator.apply_tacky_design()
            
            # Step 4: Apply content transformation
            transformer = ContentTransformer(presentation, intensity=7)
            transformer.transform_all_content()
            
            # Step 5: Save output
            output_file = os.path.join(temp_dir, "output.pptx")
            presentation.save(output_file)
            
            # Verify output exists
            assert os.path.exists(output_file)
            assert os.path.getsize(output_file) > 0


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
