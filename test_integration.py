#!/usr/bin/env python
"""Integration test: simulate the full Flask API flow"""

import os
import sys
import shutil
from pathlib import Path

# Add to path
sys.path.insert(0, str(Path(__file__).parent))

from werkzeug.utils import secure_filename
from pptx import Presentation

from src.ppt_analyzer import PPTAnalyzer
from src.taco_generator import TacoGenerator
from src.content_transformer import ContentTransformer
from config import Config

print("=== Full API flow simulation ===\n")

# Setup
os.makedirs(Config.get_upload_folder(), exist_ok=True)
test_file = 'samples/test_sample.pptx'

if not os.path.exists(test_file):
    print(f'✗ Test file not found: {test_file}')
    sys.exit(1)

print(f'Test file: {test_file}')
print(f'Upload folder: {Config.get_upload_folder()}')
print(f'Upload folder exists: {os.path.exists(Config.get_upload_folder())}')

# Simulate the exact Flask /api/process flow
print("\n--- Simulating /api/process ---\n")

# Step 1: Validate and save input file
filename = secure_filename('test_sample.pptx')
temp_input_path = os.path.join(Config.get_upload_folder(), f"input_12345.pptx")

print(f"1. Copying input file...")
shutil.copy(test_file, temp_input_path)
print(f"   ✓ Saved: {temp_input_path}")
print(f"   ✓ Size: {os.path.getsize(temp_input_path)} bytes")

# Step 2: Analyze
print(f"\n2. Analyzing...")
try:
    analyzer = PPTAnalyzer(temp_input_path)
    analysis_before = analyzer.analyze()
    print(f"   ✓ Analysis complete: {len(analysis_before.fonts)} fonts found")
except Exception as e:
    print(f"   ✗ Error: {e}")
    sys.exit(1)

# Step 3: Load and modify
print(f"\n3. Loading presentation...")
try:
    presentation = Presentation(temp_input_path)
    print(f"   ✓ Loaded: {len(presentation.slides)} slides")
except Exception as e:
    print(f"   ✗ Error: {e}")
    sys.exit(1)

# Step 4: Apply tacky design
print(f"\n4. Applying tacky design...")
try:
    design_generator = TacoGenerator(presentation, tacky_level=7, seed=42)
    design_generator.apply_tacky_design()
    print(f"   ✓ Design applied")
except Exception as e:
    print(f"   ✗ Error: {e}")
    import traceback
    traceback.print_exc()

# Step 5: Apply content transformation
print(f"\n5. Applying content transformation...")
try:
    content_transformer = ContentTransformer(presentation, intensity=7, seed=42)
    content_transformer.transform_all_content()
    print(f"   ✓ Content transformed")
except Exception as e:
    print(f"   ✗ Error: {e}")
    import traceback
    traceback.print_exc()

# Step 6: Save output (THIS IS WHERE THE ISSUE LIKELY IS)
print(f"\n6. Saving output file...")
output_filename = 'test_sample_TACKY.pptx'
temp_output_path = os.path.join(Config.get_upload_folder(), output_filename)

try:
    print(f"   Saving to: {temp_output_path}")
    presentation.save(temp_output_path)
    print(f"   ✓ Presentation.save() completed")
    
    # Verify file was created
    if not os.path.exists(temp_output_path):
        print(f"   ✗ File does not exist after save!")
        sys.exit(1)
    
    output_size = os.path.getsize(temp_output_path)
    print(f"   ✓ File exists: {output_size} bytes")
    
    if output_size == 0:
        print(f"   ✗ File is empty!")
        sys.exit(1)
    
except Exception as e:
    print(f"   ✗ Error saving: {e}")
    import traceback
    traceback.print_exc()
    sys.exit(1)

# Step 7: Verify by re-opening
print(f"\n7. Verifying downloaded file...")
try:
    # Simulate downloading: read file into memory
    with open(temp_output_path, 'rb') as f:
        file_data = f.read()
    
    print(f"   ✓ Read from disk: {len(file_data)} bytes")
    
    # Try to open from memory
    from io import BytesIO
    buffer = BytesIO(file_data)
    verify_prs = Presentation(buffer)
    print(f"   ✓ Presentation opens from memory: {len(verify_prs.slides)} slides")
    
except Exception as e:
    print(f"   ✗ Error: {e}")
    import traceback
    traceback.print_exc()
    sys.exit(1)

# Cleanup
os.remove(temp_input_path)
os.remove(temp_output_path)

print("\n--- Integration test PASSED! ---")
print("\nIf the actual Flask API is failing, the issue is likely in:")
print("1. Flask request handling")
print("2. File permissions")
print("3. Special character encoding in filenames")
