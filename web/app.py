"""Flask Web Application for DasaMaker"""

import logging
import os
import uuid
from pathlib import Path
from tempfile import TemporaryDirectory
from werkzeug.utils import secure_filename
from functools import wraps
from datetime import datetime, timedelta

from flask import Flask, render_template, request, send_file, jsonify, after_this_request
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

# Import our modules
import sys
sys.path.insert(0, str(Path(__file__).parent.parent))
from src.ppt_analyzer import PPTAnalyzer
from src.taco_generator import TacoGenerator
from src.content_transformer import ContentTransformer

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Flask app configuration
from config import Config

app = Flask(__name__)
# Apply config values
app.config['MAX_CONTENT_LENGTH'] = Config.MAX_FILE_SIZE
app.config['UPLOAD_FOLDER'] = Config.get_upload_folder()
app.config['JSON_SORT_KEYS'] = False

# Create upload folder if it doesn't exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Allowed file extension
ALLOWED_EXTENSIONS = Config.ALLOWED_EXTENSIONS

# Rate limiting tracker (in-memory, simple implementation)
request_timestamps = {}

def rate_limit(max_requests=20, window_seconds=3600):
    """Simple rate limiter decorator"""
    def decorator(f):
        @wraps(f)
        def decorated_function(*args, **kwargs):
            client_ip = request.remote_addr
            now = datetime.now()
            
            # Clean up old entries
            if client_ip in request_timestamps:
                request_timestamps[client_ip] = [
                    ts for ts in request_timestamps[client_ip]
                    if now - ts < timedelta(seconds=window_seconds)
                ]
            else:
                request_timestamps[client_ip] = []
            
            # Check rate limit
            if len(request_timestamps[client_ip]) >= max_requests:
                logger.warning(f"Rate limit exceeded for IP: {client_ip}")
                return jsonify({'error': 'リクエストが多すぎます。しばらく待ってからお試しください。'}), 429
            
            request_timestamps[client_ip].append(now)
            return f(*args, **kwargs)
        return decorated_function
    return decorator


def allowed_file(filename: str) -> bool:
    """Check if file extension is allowed"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def generate_output_filename(original_filename: str) -> str:
    """Generate output filename with TACKY suffix"""
    name, ext = os.path.splitext(original_filename)
    return f"{name}_TACKY.pptx"


def ensure_og_image():
    """Generate a social share OG image if missing (1200x630 PNG)."""
    try:
        static_dir = app.static_folder or str(Path(__file__).parent / 'static')
        og_path = Path(static_dir) / 'og-image.png'
        # If exists and size > 0, skip
        if og_path.exists() and og_path.stat().st_size > 0:
            return

        # Create image
        width, height = 1200, 630
        img = Image.new('RGB', (width, height), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Tacky neon stripes background
        stripes = ['#ff00ff', '#00ffff', '#ffff00', '#ff0000', '#00ff00', '#ff8800']
        stripe_h = max(1, height // len(stripes))
        for i, color in enumerate(stripes):
            y0 = i * stripe_h
            y1 = (i + 1) * stripe_h if i < len(stripes) - 1 else height
            draw.rectangle([0, y0, width, y1], fill=color)

        # White panel to simulate a slide
        panel_margin = 60
        draw.rounded_rectangle([panel_margin, panel_margin, width - panel_margin, height - 180], radius=24, fill='#ffffff', outline='#00ff00', width=8)

        # Try loading a common bold font (Windows)
        font_paths = [
            'C:/Windows/Fonts/arialbd.ttf',
            'C:/Windows/Fonts/ARIALBD.TTF',
            'C:/Windows/Fonts/impact.ttf',
            'C:/Windows/Fonts/IMPACT.TTF'
        ]
        title_font = None
        subtitle_font = None
        for p in font_paths:
            try:
                title_font = ImageFont.truetype(p, 120)
                subtitle_font = ImageFont.truetype(p, 48)
                break
            except Exception:
                continue
        if title_font is None:
            title_font = ImageFont.load_default()
        if subtitle_font is None:
            subtitle_font = ImageFont.load_default()

        # Title text
        title = 'UglySlide'
        tw, th = draw.textlength(title, font=title_font), 120
        tx = (width - tw) / 2
        ty = height - 160
        # Outline effect
        for dx, dy in [(-3,0),(3,0),(0,-3),(0,3)]:
            draw.text((tx+dx, ty+dy), title, font=title_font, fill='#000000')
        draw.text((tx, ty), title, font=title_font, fill='#ffffff')

        # Subtitle
        subtitle = 'PowerPoint Uncooler — Make it intentionally tacky!'
        sw = draw.textlength(subtitle, font=subtitle_font)
        sx = (width - sw) / 2
        sy = height - 80
        draw.text((sx, sy), subtitle, font=subtitle_font, fill='#000000')

        # Save PNG
        og_path.parent.mkdir(parents=True, exist_ok=True)
        img.save(str(og_path), format='PNG')
    except Exception as e:
        logger.warning(f"Failed to generate OG image: {e}")


@app.route('/')
def index():
    """Render main page"""
    # Ensure OG image exists for social share
    ensure_og_image()
    return render_template('index.html')


@app.route('/about')
def about():
    """Render about page"""
    return render_template('about.html')


@app.route('/privacy')
def privacy():
    """Render privacy policy page"""
    return render_template('privacy.html')


@app.route('/terms')
def terms():
    """Render terms of service page"""
    return render_template('terms.html')


@app.route('/favicon.ico')
def favicon():
    """Serve favicon for browsers requesting /favicon.ico"""
    try:
        favicon_path = os.path.join(app.static_folder or 'static', 'favicon.svg')
        if os.path.exists(favicon_path):
            return send_file(favicon_path, mimetype='image/svg+xml')
        # If missing, return 404 gracefully
        return jsonify({'error': 'favicon not found'}), 404
    except Exception:
        return jsonify({'error': 'failed to serve favicon'}), 500


@app.route('/api/process', methods=['POST'])
@rate_limit(max_requests=10, window_seconds=3600)
def process_presentation():
    """
    Process uploaded PPTX file and generate tacky version
    
    Query parameters:
    - design_level: Design tackiness level (1-10, default=7)
    - content_level: Content transformation intensity (1-10, default=7)
    - seed: Optional integer seed for deterministic output
    """
    
    try:
        # Check if file is present
        if 'file' not in request.files:
            return jsonify({'error': 'ファイルが選択されていません'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'ファイルが選択されていません'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'error': 'PPTXファイルのみ対応しています'}), 400
        
        # Validate filename for path traversal
        try:
            secure_input_filename = secure_filename(file.filename)
            if not secure_input_filename or secure_input_filename == '':
                return jsonify({'error': 'ファイル名が無効です'}), 400
        except Exception:
            return jsonify({'error': 'ファイル名が無効です'}), 400
        
        # Get tackiness levels from request
        try:
            design_level = int(request.args.get('design_level', 7))
            content_level = int(request.args.get('content_level', 7))
            
            # Validate ranges
            if not (1 <= design_level <= 10):
                design_level = 7
            if not (1 <= content_level <= 10):
                content_level = 7
                
            seed = request.args.get('seed')
            seed = int(seed) if seed is not None and str(seed).isdigit() else None
        except (ValueError, TypeError):
            return jsonify({'error': 'パラメータが無効です'}), 400
        
        # Create temporary file
        temp_input_path = os.path.join(
            app.config['UPLOAD_FOLDER'],
            f"input_{uuid.uuid4().hex}.pptx"
        )
        
        # Save uploaded file
        try:
            file.save(temp_input_path)
            logger.info(f"File saved: {temp_input_path}")
        except Exception as e:
            logger.error(f"Failed to save file: {e}")
            return jsonify({'error': 'ファイルの保存に失敗しました'}), 500
        
        try:
            # Check file size again after save (just in case)
            file_size = os.path.getsize(temp_input_path)
            if file_size > Config.MAX_FILE_SIZE:
                return jsonify({'error': 'ファイルサイズが制限を超えています'}), 413
            
            # Step 1: Analyze original presentation
            logger.info("Step 1: Analyzing presentation...")
            analyzer = PPTAnalyzer(temp_input_path)
            analysis_before = analyzer.analyze()
            logger.info(f"Analysis complete: {analysis_before.to_dict()}")
            
            # Step 2: Load presentation for modifications
            logger.info("Step 2: Loading presentation for modifications...")
            presentation = Presentation(temp_input_path)
            
            # Step 3: Apply tacky design
            logger.info(f"Step 3: Applying tacky design (level {design_level})...")
            design_generator = TacoGenerator(presentation, tacky_level=design_level, seed=seed)
            design_generator.apply_tacky_design()
            
            # Step 4: Apply content transformation
            logger.info(f"Step 4: Applying content transformation (intensity {content_level})...")
            content_transformer = ContentTransformer(presentation, intensity=content_level, seed=seed)
            content_transformer.transform_all_content()
            
            # Step 5: Save output
            output_filename = generate_output_filename(secure_filename(file.filename))
            temp_output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
            
            logger.info(f"Step 5: Saving output to {temp_output_path}...")
            try:
                presentation.save(temp_output_path)
                # Verify file was created and has content
                if not os.path.exists(temp_output_path):
                    raise Exception(f"Output file was not created at {temp_output_path}")
                file_size = os.path.getsize(temp_output_path)
                if file_size == 0:
                    raise Exception(f"Output file is empty (0 bytes)")
                logger.info(f"Output file saved successfully: {temp_output_path} ({file_size} bytes)")
                
                # Additional verification: try to open the file to ensure it's not corrupted
                logger.info("Verifying saved file integrity...")
                try:
                    verify_prs = Presentation(temp_output_path)
                    slide_count = len(verify_prs.slides)
                    logger.info(f"✓ File integrity verified: {slide_count} slides, {file_size} bytes")
                except Exception as verify_e:
                    logger.error(f"File integrity check failed: {verify_e}")
                    raise Exception(f"Saved file is corrupted or unreadable: {verify_e}")
                    
            except Exception as e:
                logger.error(f"Failed to save presentation: {e}", exc_info=True)
                # Clean up corrupted file
                if os.path.exists(temp_output_path):
                    try:
                        os.remove(temp_output_path)
                    except:
                        pass
                raise
            
            # Optionally analyze after modifications for before/after comparison
            try:
                analysis_after = PPTAnalyzer(temp_output_path).analyze()
            except Exception as e:
                logger.warning(f"Failed to analyze output file: {e}")
                analysis_after = None

            logger.info("Processing complete!")
            
            # Return success with file info
            response_payload = {
                'success': True,
                'filename': output_filename,
                'seed': seed,
                'analysis': {
                    'total_slides': analysis_before.total_slides,
                    'fonts_found': len(analysis_before.fonts),
                    'colors_found': len(analysis_before.colors),
                    'animations_found': analysis_before.animation_count,
                }
            }
            if analysis_after is not None:
                response_payload['analysis_before'] = response_payload.pop('analysis')
                response_payload['analysis_after'] = {
                    'total_slides': analysis_after.total_slides,
                    'fonts_found': len(analysis_after.fonts),
                    'colors_found': len(analysis_after.colors),
                    'animations_found': analysis_after.animation_count,
                }

            return jsonify(response_payload)
        
        finally:
            # Clean up input file
            if os.path.exists(temp_input_path):
                try:
                    os.remove(temp_input_path)
                    logger.info(f"Cleaned up input file: {temp_input_path}")
                except Exception as e:
                    logger.warning(f"Failed to clean up input file: {e}")
    
    except Exception as e:
        logger.error(f"Error processing file: {e}", exc_info=True)
        return jsonify({'error': 'ファイルの処理に失敗しました。ファイル形式が正しいか確認してください。'}), 500


@app.route('/api/download/<filename>')
def download_file(filename: str):
    """
    Download processed PPTX file
    
    Args:
        filename: Name of the file to download
    """
    
    try:
        # Validate filename (prevent directory traversal)
        filename = secure_filename(filename)
        if not filename:
            return jsonify({'error': 'ファイル名が無効です'}), 400
            
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        # Check if file exists
        if not os.path.exists(filepath):
            logger.warning(f"Download requested for non-existent file: {filename}")
            return jsonify({'error': 'ファイルが見つかりません'}), 404
        
        # Check if file is in the correct directory (path traversal prevention)
        if not os.path.abspath(filepath).startswith(os.path.abspath(app.config['UPLOAD_FOLDER'])):
            logger.error(f"Potential path traversal attempt: {filename}")
            return jsonify({'error': 'ファイルアクセスが拒否されました'}), 403
        
        # Check file size
        file_size = os.path.getsize(filepath)
        if file_size == 0:
            logger.warning(f"Zero-sized file detected: {filename}")
            return jsonify({'error': 'ファイルが壊れています'}), 400
        
        logger.info(f"Downloading file: {filename} (size: {file_size} bytes)")

        @after_this_request
        def remove_file(response):
            try:
                if os.path.exists(filepath):
                    os.remove(filepath)
                    logger.info(f"Cleaned up file after send: {filepath}")
            except Exception as e:
                logger.debug(f"Could not clean up file after send: {e}")
            return response

        # Use filepath directly with send_file for more reliable delivery
        response = send_file(
            filepath,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )
        response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
        response.headers['Pragma'] = 'no-cache'
        response.headers['Expires'] = '0'
        logger.info(f"File sent successfully: {filename}")
        return response
    
    except Exception as e:
        logger.error(f"Error downloading file: {e}", exc_info=True)
        return jsonify({'error': 'ダウンロードに失敗しました'}), 500


@app.route('/api/health')
def health_check():
    """Health check endpoint"""
    return jsonify({
        'status': 'ok',
        'version': '0.1.0',
        'timestamp': datetime.now().isoformat()
    })


@app.errorhandler(413)
def request_entity_too_large(error):
    """Handle file too large error"""
    logger.warning(f"File too large error: {error}")
    return jsonify({'error': 'ファイルサイズが大きすぎます（最大50MB）'}), 413


@app.errorhandler(404)
def not_found(error):
    """Handle 404 errors"""
    logger.warning(f"404 error: {error}")
    return jsonify({'error': 'ページが見つかりません'}), 404


@app.errorhandler(500)
def internal_error(error):
    """Handle 500 errors"""
    logger.error(f"Internal server error: {error}")
    return jsonify({'error': 'サーバーエラーが発生しました。管理者に連絡してください。'}), 500


@app.errorhandler(429)
def rate_limited(error):
    """Handle rate limit errors"""
    logger.warning(f"Rate limit exceeded: {error}")
    return jsonify({'error': 'リクエストが多すぎます。しばらく待ってからお試しください。'}), 429


# Security headers
@app.after_request
def set_security_headers(response):
    """Add security headers to responses"""
    response.headers['X-Content-Type-Options'] = 'nosniff'
    response.headers['X-Frame-Options'] = 'SAMEORIGIN'
    response.headers['X-XSS-Protection'] = '1; mode=block'
    response.headers['Strict-Transport-Security'] = 'max-age=31536000; includeSubDomains'
    return response


if __name__ == '__main__':
    # Check environment
    debug_mode = os.getenv('FLASK_ENV', 'production') == 'development'
    
    # Get port from environment variable (Render.com uses PORT env var)
    port = int(os.getenv('PORT', 5000))
    
    logger.info("=" * 60)
    logger.info("Starting DasaMaker Flask application...")
    logger.info(f"Debug mode: {debug_mode}")
    logger.info(f"Upload folder: {app.config['UPLOAD_FOLDER']}")
    logger.info(f"Max file size: {app.config['MAX_CONTENT_LENGTH'] / (1024*1024):.1f} MB")
    logger.info(f"Port: {port}")
    logger.info("=" * 60)
    
    app.run(debug=debug_mode, host='0.0.0.0', port=port, threaded=True)

