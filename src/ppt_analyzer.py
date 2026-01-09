"""PPT Analyzer module - Extract design elements from PPTX files"""

import logging
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Any

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


class DesignAnalysis:
    """Store extracted design elements from a presentation"""
    
    def __init__(self):
        self.fonts: Dict[str, int] = {}  # Font name -> count
        self.colors: List[Tuple[int, int, int]] = []  # RGB tuples
        self.slide_layouts: List[str] = []
        self.animation_count = 0
        self.total_slides = 0
        self.text_elements: List[str] = []
        
    def to_dict(self) -> Dict[str, Any]:
        """Convert analysis to dictionary"""
        return {
            'fonts': self.fonts,
            'colors': self.colors,
            'slide_layouts': self.slide_layouts,
            'animation_count': self.animation_count,
            'total_slides': self.total_slides,
            'text_elements_count': len(self.text_elements),
        }


class PPTAnalyzer:
    """Analyze PowerPoint presentations to extract design elements"""
    
    def __init__(self, pptx_path: str):
        """
        Initialize analyzer with a PPTX file
        
        Args:
            pptx_path: Path to the PPTX file
            
        Raises:
            FileNotFoundError: If file does not exist
            ValueError: If file is not a valid PPTX
        """
        self.pptx_path = Path(pptx_path)
        
        if not self.pptx_path.exists():
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
        
        if not self.pptx_path.suffix.lower() == '.pptx':
            raise ValueError(f"File must be PPTX format: {pptx_path}")
        
        try:
            self.presentation = Presentation(str(self.pptx_path))
        except Exception as e:
            raise ValueError(f"Failed to parse PPTX: {e}")
    
    def analyze(self) -> DesignAnalysis:
        """
        Analyze the entire presentation
        
        Returns:
            DesignAnalysis object with extracted design elements
        """
        analysis = DesignAnalysis()
        analysis.total_slides = len(self.presentation.slides)
        
        logger.info(f"Analyzing presentation: {self.pptx_path.name}")
        logger.info(f"Total slides: {analysis.total_slides}")
        
        # Analyze each slide
        for slide_idx, slide in enumerate(self.presentation.slides):
            self._analyze_slide(slide, slide_idx, analysis)
        
        logger.info(f"Design analysis complete: {analysis.to_dict()}")
        return analysis
    
    def _analyze_slide(self, slide, slide_idx: int, analysis: DesignAnalysis) -> None:
        """Analyze a single slide"""
        
        # Store layout name
        try:
            layout_name = slide.slide_layout.name
            if layout_name not in analysis.slide_layouts:
                analysis.slide_layouts.append(layout_name)
        except Exception as e:
            logger.debug(f"Could not extract layout from slide {slide_idx}: {e}")
        
        # Analyze shapes
        for shape in slide.shapes:
            self._analyze_shape(shape, analysis)
        
        # Count animations/transitions (approximation)
        try:
            # python-pptx does not expose full animation API; count transitions as proxy
            trans = getattr(slide, 'slide_show_transition', None)
            if trans is not None:
                # Heuristic: count a transition if any notable attribute is set
                has_effect = any([
                    getattr(trans, 'entry_effect', None),
                    getattr(trans, 'advance_on_click', None),
                    getattr(trans, 'advance_after_time', None),
                ])
                if has_effect:
                    analysis.animation_count += 1
        except Exception as e:
            logger.debug(f"Could not extract transitions from slide {slide_idx}: {e}")
    
    def _analyze_shape(self, shape, analysis: DesignAnalysis) -> None:
        """Analyze a single shape for fonts and colors"""
        
        # Extract text
        if hasattr(shape, "text_frame"):
            try:
                text = shape.text.strip()
                if text:
                    analysis.text_elements.append(text)
                
                # Extract font information
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            font_name = run.font.name
                            analysis.fonts[font_name] = analysis.fonts.get(font_name, 0) + 1
                        
                        # Extract color
                        if run.font.color.type == 1:  # RGB color
                            try:
                                rgb = run.font.color.rgb
                                rgb_tuple = (rgb[0], rgb[1], rgb[2])
                                if rgb_tuple not in analysis.colors:
                                    analysis.colors.append(rgb_tuple)
                            except Exception:
                                pass
            except Exception as e:
                logger.debug(f"Could not analyze text in shape: {e}")
        
        # Extract fill color
        if hasattr(shape, "fill"):
            try:
                fill = shape.fill
                if fill.type == 1:  # SOLID fill
                    try:
                        rgb = fill.fore_color.rgb
                        rgb_tuple = (rgb[0], rgb[1], rgb[2])
                        if rgb_tuple not in analysis.colors:
                            analysis.colors.append(rgb_tuple)
                    except Exception:
                        pass
            except Exception as e:
                logger.debug(f"Could not analyze fill color: {e}")
    
    def get_dominant_fonts(self, top_n: int = 5) -> List[Tuple[str, int]]:
        """
        Get most common fonts in the presentation
        
        Args:
            top_n: Number of top fonts to return
            
        Returns:
            List of (font_name, count) tuples sorted by frequency
        """
        analysis = self.analyze()
        sorted_fonts = sorted(analysis.fonts.items(), key=lambda x: x[1], reverse=True)
        return sorted_fonts[:top_n]
    
    def get_color_palette(self) -> List[Tuple[int, int, int]]:
        """
        Get color palette from the presentation
        
        Returns:
            List of RGB color tuples
        """
        analysis = self.analyze()
        return analysis.colors[:10]  # Limit to top 10 colors


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    # Example usage
    # analyzer = PPTAnalyzer("sample.pptx")
    # analysis = analyzer.analyze()
    # print(analysis.to_dict())
