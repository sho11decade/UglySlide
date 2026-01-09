"""Content Transformer module - Apply satirical content modifications"""

import logging
import random
from typing import List, Optional

from pptx import Presentation

logger = logging.getLogger(__name__)


# Verbose filler phrases
VERBOSE_FILLERS = [
    "In today's dynamic and rapidly evolving landscape,",
    "It is imperative to note that,",
    "As we navigate the complexities of modern business,",
    "Synergistically speaking,",
    "With all due respect to the aforementioned points,",
    "In the spirit of full transparency,",
    "One could argue that,",
    "From a holistic perspective,",
]

# Outdated corporate jargon
OUTDATED_JARGON = [
    "paradigm shift",
    "leverage synergies",
    "circle back",
    "deep dive",
    "move the needle",
    "boil the ocean",
    "touch base",
    "take it offline",
    "think outside the box",
    "win-win situation",
]

# Random tangential facts to insert
TANGENTIAL_FACTS = [
    " (Did you know? Bananas are berries, but strawberries are not!)",
    " (Fun fact: A group of flamingos is called a 'flamboyance'!)",
    " (Interesting: The shortest war in history lasted 38 minutes!)",
    " (Trivia: Honey never spoils and can last forever!)",
    " (Fun fact: Octopuses have three hearts!)",
    " (Random fact: The Great Wall of China is NOT visible from space!)",
]

# Tone modifications
SARCASM_PREFIXES = [
    "Let me tell you, ",
    "Buckle up, because ",
    "Oh, how exciting - ",
    "Shockingly, ",
    "You won't believe this, but ",
]


class ContentTransformer:
    """Transform presentation content with satirical modifications"""
    
    def __init__(self, presentation: Presentation, intensity: int = 7, seed: Optional[int] = None):
        """
        Initialize content transformer
        
        Args:
            presentation: python-pptx Presentation object
            intensity: How aggressive the transformation is (1-10), default=7
        """
        self.presentation = presentation
        self.intensity = max(1, min(10, intensity))  # Clamp to 1-10
        # Deterministic RNG if seed provided
        self._rand = random.Random(seed) if seed is not None else random.Random()
        logger.info(f"ContentTransformer initialized with intensity: {self.intensity}")
    
    def transform_all_content(self) -> None:
        """Apply satirical transformations to all slides"""
        logger.info(f"Starting content transformation (intensity {self.intensity})")
        
        for slide_idx, slide in enumerate(self.presentation.slides):
            logger.info(f"Transforming slide {slide_idx + 1}/{len(self.presentation.slides)}")
            self._transform_slide_content(slide)
    
    def _transform_slide_content(self, slide) -> None:
        """Apply transformations to a single slide"""
        
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                self._transform_text_shape(shape)
    
    def _transform_text_shape(self, shape) -> None:
        """Transform text in a shape with satirical modifications"""
        
        try:
            text_frame = shape.text_frame
            
            for paragraph in text_frame.paragraphs:
                # Get original text
                original_text = paragraph.text.strip()
                
                if not original_text:
                    continue
                
                # Apply various transformations based on intensity
                modified_text = original_text
                
                # Level 1-3: Add verbose filler at start
                if self.intensity >= 1 and self._rand.random() < 0.3 * (self.intensity / 10):
                    modified_text = f"{self._rand.choice(VERBOSE_FILLERS)} {modified_text}"
                
                # Level 4-6: Insert outdated jargon
                if self.intensity >= 4 and self._rand.random() < 0.4 * (self.intensity / 10):
                    jargon = self._rand.choice(OUTDATED_JARGON)
                    # Insert jargon naturally into text
                    words = modified_text.split()
                    if len(words) > 2:
                        insert_pos = self._rand.randint(1, len(words) - 1)
                        words.insert(insert_pos, f"({jargon})")
                        modified_text = " ".join(words)
                
                # Level 5-7: Add tangential facts
                if self.intensity >= 5 and self._rand.random() < 0.25 * (self.intensity / 10):
                    modified_text += self._rand.choice(TANGENTIAL_FACTS)
                
                # Level 7-9: Add sarcasm/tone changes
                if self.intensity >= 7 and self._rand.random() < 0.2:
                    modified_text = f"{self._rand.choice(SARCASM_PREFIXES)}{modified_text}"
                
                # Level 9-10: Repeat key phrases for emphasis
                if self.intensity >= 9 and self._rand.random() < 0.15:
                    words = modified_text.split()
                    if len(words) > 0:
                        key_word = self._rand.choice(words[-3:])
                        modified_text += f" Did we mention {key_word}? Yes, {key_word} is very important."
                
                # Replace paragraph text
                if modified_text != original_text:
                    # Clear existing runs
                    for run in paragraph.runs:
                        run.text = ""
                    
                    # Add new text
                    if len(paragraph.runs) > 0:
                        paragraph.runs[0].text = modified_text
                    else:
                        paragraph.text = modified_text
                    
                    logger.debug(f"Transformed: '{original_text}' → '{modified_text}'")
        
        except Exception as e:
            logger.debug(f"Could not transform text shape: {e}")
    
    def get_intensity(self) -> int:
        """Get current transformation intensity (1-10)"""
        return self.intensity
    
    def set_intensity(self, intensity: int) -> None:
        """Set transformation intensity (1-10)"""
        self.intensity = max(1, min(10, intensity))
        logger.info(f"Transformation intensity set to: {self.intensity}")


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    # Example usage
    # presentation = Presentation("sample.pptx")
    # transformer = ContentTransformer(presentation, intensity=7)
    # transformer.transform_all_content()
    # presentation.save("sample_CONTENT_TRANSFORMED.pptx")
