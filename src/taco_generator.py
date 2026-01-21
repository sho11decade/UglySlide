"""Taco Generator module - Generate "uncool" design variations"""

import logging
import random
from typing import List, Tuple, Optional, Dict, Any

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from tempfile import NamedTemporaryFile
from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger(__name__)


# Tacky font replacement strategy
TACKY_FONTS = [
    "Comic Sans MS",
    "Papyrus",
    "Impact",
    "Curlz MT",
    "Jokerman",
    "Brush Script MT",
    "Chiller",
    "Lucida Handwriting",
]

# Neon color palette for tacky transformations
NEON_COLORS = [
    (255, 0, 127),    # Hot pink
    (0, 255, 255),    # Cyan
    (255, 255, 0),    # Bright yellow
    (0, 255, 0),      # Lime green
    (255, 127, 0),    # Bright orange
    (255, 0, 0),      # Bright red
    (148, 0, 211),    # Blue-violet
    (255, 192, 203),  # Light pink
]

# Extreme neon color combinations for gradients
EXTREME_NEON_PAIRS = [
    [(255, 0, 127), (0, 255, 255)],      # Hot pink to cyan
    [(255, 255, 0), (255, 0, 0)],        # Yellow to red
    [(0, 255, 0), (148, 0, 211)],        # Lime to blue-violet
    [(255, 127, 0), (255, 0, 127)],      # Orange to pink
    [(0, 255, 255), (255, 255, 0)],      # Cyan to yellow
    [(148, 0, 211), (255, 0, 0)],        # Blue-violet to red
    [(255, 192, 203), (0, 255, 0)],      # Light pink to lime
    [(255, 0, 0), (0, 255, 255)],        # Red to cyan
    [(255, 255, 0), (148, 0, 211)],      # Yellow to blue-violet
    [(255, 127, 0), (0, 255, 0)],        # Orange to lime
]


class TacoGenerator:
    """Generate intentionally uncool design variations of PowerPoint presentations"""
    
    def __init__(self, presentation: Presentation, tacky_level: int = 7, seed: Optional[int] = None):
        """
        Initialize tacky design generator
        
        Args:
            presentation: python-pptx Presentation object
            tacky_level: Intensity of tackiness (1-10), default=7
        """
        self.presentation = presentation
        self.tacky_level = max(1, min(10, tacky_level))  # Clamp to 1-10
        # Deterministic RNG if seed provided
        self._rand = random.Random(seed) if seed is not None else random.Random()
        logger.info(f"TacoGenerator initialized with tackiness level: {self.tacky_level}")
    
    def apply_tacky_design(self) -> None:
        """Apply all tacky design transformations to the presentation"""
        logger.info(f"Applying tacky design transformations (level {self.tacky_level})")
        
        for slide_idx, slide in enumerate(self.presentation.slides):
            logger.info(f"Processing slide {slide_idx + 1}/{len(self.presentation.slides)}")
            self._make_slide_tacky(slide)
    
    def _make_slide_tacky(self, slide) -> None:
        """Apply tacky transformations to a single slide"""
        
        # Apply tacky fonts and colors to text
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                self._tacky_text_transform(shape)
            
            # Apply tacky fill colors to shapes
            if hasattr(shape, "fill") and shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                self._tacky_fill_transform(shape)
                
                # Apply extreme gradient transformations
                if self.tacky_level >= 5:
                    self._apply_extreme_gradient(shape)
        
        # Apply background color if tackiness is high
        if self.tacky_level >= 6:
            self._apply_tacky_background(slide)
            # Transition insertion disabled due to corruption reports
            # self._apply_overkill_transition(slide)

        # Add gaudy footer banner
        if self.tacky_level >= 7:
            self._add_footer_banner(slide)

        # Randomize positions and rotations for extra chaos
        if self.tacky_level >= 8:
            self._randomize_layout(slide)
            # Add some stickers
            self._insert_gaudy_stickers(slide, count=self._rand.randint(1, 3))

        # Set a picture background and drop in a chart at extreme levels
        if self.tacky_level >= 9:
            self._set_background_picture(slide)
            self._insert_gaudy_chart(slide)
    
    def _tacky_text_transform(self, shape) -> None:
        """Transform text with tacky fonts and colors"""
        
        try:
            text_frame = shape.text_frame
            
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    # Replace with tacky font
                    if self.tacky_level >= 3:
                        run.font.name = self._rand.choice(TACKY_FONTS)
                    
                    # Apply tacky color
                    if self.tacky_level >= 5:
                        run.font.color.rgb = RGBColor(*self._rand.choice(NEON_COLORS))
                    
                    # Increase font size for emphasis
                    if self.tacky_level >= 7 and run.font.size:
                        current_size = run.font.size.pt
                        run.font.size = Pt(current_size * 1.3)
                    
                    # Make everything bold and italic for extra tackiness
                    if self.tacky_level >= 8:
                        run.font.bold = True
                        run.font.italic = True
        except Exception as e:
            logger.debug(f"Could not transform text in shape: {e}")
    
    def _tacky_fill_transform(self, shape) -> None:
        """Transform shape fill with tacky colors"""
        
        try:
            if self.tacky_level >= 6:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*self._rand.choice(NEON_COLORS))
        except Exception as e:
            logger.debug(f"Could not transform shape fill: {e}")
    
    def _apply_tacky_background(self, slide) -> None:
        """Apply tacky background color to slide"""
        
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            # Use a less intense color for background
            fill.fore_color.rgb = RGBColor(200 + self._rand.randint(-50, 50), 
                                          100 + self._rand.randint(-50, 50),
                                          150 + self._rand.randint(-50, 50))
        except Exception as e:
            logger.debug(f"Could not apply background: {e}")
    
    def _apply_extreme_gradient(self, shape) -> None:
        """Apply extreme gradient fills to shapes with leveled intensity"""
        
        try:
            fill = shape.fill
            
            # Level 5-6: Simple two-color gradients
            if self.tacky_level in [5, 6]:
                self._apply_simple_gradient(fill)
            
            # Level 7-8: More intense multi-color gradients
            elif self.tacky_level in [7, 8]:
                self._apply_multi_color_gradient(fill)
            
            # Level 9-10: Absolutely crazy extreme gradients
            elif self.tacky_level >= 9:
                self._apply_extreme_multi_gradient(fill)
                
            logger.debug(f"Applied gradient (level {self.tacky_level}) to shape")
        except Exception as e:
            logger.debug(f"Could not apply gradient: {e}")
    
    def _apply_simple_gradient(self, fill) -> None:
        """Apply simple two-color gradient (level 5-6)"""
        
        try:
            fill.gradient()
            fill.gradient_angle = self._rand.choice([0, 45, 90, 135, 180, 225, 270, 315])
            
            # Get random extreme color pair
            color_pair = self._rand.choice(EXTREME_NEON_PAIRS)
            
            fill.gradient_stops[0].color.rgb = RGBColor(*color_pair[0])
            fill.gradient_stops[1].color.rgb = RGBColor(*color_pair[1])
            
            logger.debug(f"Applied simple gradient: {color_pair[0]} to {color_pair[1]}")
        except Exception as e:
            logger.debug(f"Could not apply simple gradient: {e}")
    
    def _apply_multi_color_gradient(self, fill) -> None:
        """Apply multi-color gradient with 3-4 color stops (level 7-8)"""
        
        try:
            fill.gradient()
            fill.gradient_angle = self._rand.choice([0, 45, 90, 135, 180, 225, 270, 315])
            
            # Select 3-4 neon colors
            num_colors = self._rand.choice([3, 4])
            colors = self._rand.sample(NEON_COLORS, min(num_colors, len(NEON_COLORS)))
            
            # Clear existing stops and add new ones
            while len(fill.gradient_stops) > 2:
                # Remove extra stops (keep first 2)
                try:
                    fill._element.remove(fill.gradient_stops[2]._element)
                except:
                    break
            
            # Set first two stops
            fill.gradient_stops[0].color.rgb = RGBColor(*colors[0])
            fill.gradient_stops[1].color.rgb = RGBColor(*colors[-1])
            
            # Try to add intermediate stops
            try:
                for i in range(1, len(colors) - 1):
                    # Add gradient stop at position
                    position = i / (len(colors) - 1)
                    try:
                        new_stop = fill.gradient_stops._insert_stop(position)
                        new_stop.color.rgb = RGBColor(*colors[i])
                    except:
                        pass
            except Exception as e:
                logger.debug(f"Could not add intermediate gradient stops: {e}")
            
            logger.debug(f"Applied multi-color gradient with {len(colors)} colors")
        except Exception as e:
            logger.debug(f"Could not apply multi-color gradient: {e}")
    
    def _apply_extreme_multi_gradient(self, fill) -> None:
        """Apply absolutely extreme multi-color gradient (level 9-10)"""
        
        try:
            fill.gradient()
            
            # Randomize angle
            fill.gradient_angle = self._rand.choice([0, 45, 90, 135, 180, 225, 270, 315])
            
            # Use all neon colors in extreme combinations
            colors = NEON_COLORS.copy()
            self._rand.shuffle(colors)
            
            # Set gradient stops
            fill.gradient_stops[0].color.rgb = RGBColor(*colors[0])
            fill.gradient_stops[1].color.rgb = RGBColor(*colors[-1])
            
            # Try to add even more gradient stops for maximum effect
            try:
                for i in range(1, len(colors) - 1):
                    position = i / (len(colors) - 1)
                    try:
                        new_stop = fill.gradient_stops._insert_stop(position)
                        new_stop.color.rgb = RGBColor(*colors[i])
                    except:
                        pass
            except:
                pass
            
            # Extra extreme: apply multiple gradient angles
            if self.tacky_level == 10:
                try:
                    fill.gradient_angle = self._rand.randint(0, 360)
                except:
                    pass
            
            logger.debug(f"Applied extreme multi-gradient with {len(colors)} colors")
        except Exception as e:
            logger.debug(f"Could not apply extreme gradient: {e}")

    def _apply_overkill_transition(self, slide) -> None:
        """Set a gaudy slide transition to simulate excessive animation."""

        try:
            sld = slide._element

            # Remove any existing transition
            for node in list(sld):
                if node.tag == qn('p:transition'):
                    sld.remove(node)

            transition = OxmlElement('p:transition')
            transition.set('spd', 'fast')
            transition.set('advClick', '1')
            transition.set('advTm', str(self._rand.randint(600, 1800)))

            choice = self._rand.choice([
                'fade', 'dissolve', 'checker', 'circle', 'random', 'zoom',
                'blinds', 'wipe', 'push', 'cover', 'split', 'wheel'
            ])

            try:
                if choice in {'blinds', 'wipe', 'push', 'cover'}:
                    child = OxmlElement(f'p:{choice}')
                    child.set('dir', self._rand.choice(['l', 'r', 'u', 'd']))
                elif choice == 'split':
                    child = OxmlElement('p:split')
                    child.set('dir', self._rand.choice(['l', 'r', 'u', 'd']))
                    child.set('orient', self._rand.choice(['horz', 'vert']))
                elif choice == 'wheel':
                    child = OxmlElement('p:wheel')
                    child.set('spokes', str(self._rand.choice([4, 6, 8])))
                else:
                    child = OxmlElement(f'p:{choice}')

                transition.append(child)
                sld.insert(0, transition)
            except Exception as inner_e:
                logger.debug(f"Failed to create transition element {choice}: {inner_e}")
        except Exception as e:
            logger.debug(f"Could not apply overkill transition: {e}")
    
    
    def get_tacky_level(self) -> int:
        """Get current tackiness level (1-10)"""
        return self.tacky_level
    
    def set_tacky_level(self, level: int) -> None:
        """Set tackiness level (1-10)"""
        self.tacky_level = max(1, min(10, level))
        logger.info(f"Tackiness level set to: {self.tacky_level}")

    # --- Advanced tacky helpers ---

    def _add_footer_banner(self, slide) -> None:
        """Add a neon footer banner with text across the bottom."""
        try:
            slide_width = self.presentation.slide_width
            slide_height = self.presentation.slide_height
            margin = int(slide_height * 0.06)
            height = int(slide_height * 0.10)
            left = 0
            top = slide_height - height - margin // 2
            width = slide_width

            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*self._rand.choice(NEON_COLORS))
            shape.line.width = Pt(6)
            shape.line.color.rgb = RGBColor(*self._rand.choice(NEON_COLORS))

            # Text content
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "★ UGLYSLIDE — PowerPoint Uncooler ★"
            if p.runs:
                r = p.runs[0]
                r.font.name = self._rand.choice(TACKY_FONTS)
                r.font.size = Pt(24)
                r.font.bold = True
                r.font.color.rgb = RGBColor(0, 0, 0)
        except Exception as e:
            logger.debug(f"Could not add footer banner: {e}")

    def _randomize_layout(self, slide) -> None:
        """Randomly reposition and rotate non-picture shapes."""
        try:
            sw = self.presentation.slide_width
            sh = self.presentation.slide_height
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    continue
                try:
                    w = getattr(shape, 'width', int(sw * 0.2))
                    h = getattr(shape, 'height', int(sh * 0.1))
                    new_left = self._rand.randint(0, max(0, sw - w))
                    new_top = self._rand.randint(0, max(0, sh - h))
                    shape.left = new_left
                    shape.top = new_top
                    shape.rotation = self._rand.randint(-25, 25)
                except Exception:
                    continue
        except Exception as e:
            logger.debug(f"Could not randomize layout: {e}")

    def _generate_sticker_png(self, size: int = 160) -> bytes:
        """Create a small neon sticker PNG in memory and return bytes."""
        img = Image.new('RGBA', (size, size), (0, 0, 0, 0))
        d = ImageDraw.Draw(img)
        # Starburst
        center = (size//2, size//2)
        for r in range(size//2, 10, -10):
            d.ellipse([(center[0]-r, center[1]-r), (center[0]+r, center[1]+r)], outline=self._rand.choice(['#ff00ff','#00ffff','#ffff00','#ff0000','#00ff00']), width=6)
        # Text
        txt = self._rand.choice(["WOW","LOL","★","NEON","BOOM"])
        try:
            font = ImageFont.truetype('C:/Windows/Fonts/impact.ttf', int(size*0.28))
        except Exception:
            font = ImageFont.load_default()
        tw = d.textlength(txt, font=font)
        d.text(((size-tw)/2, size*0.35), txt, font=font, fill='#000000')
        # Save to bytes
        with NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            img.save(tmp.name, format='PNG')
            tmp.seek(0)
            data = tmp.read()
            path = tmp.name
        return data, path

    def _insert_gaudy_stickers(self, slide, count: int = 2) -> None:
        """Insert a few neon sticker pictures at random positions."""
        try:
            sw = self.presentation.slide_width
            sh = self.presentation.slide_height
            for _ in range(count):
                try:
                    _data, path = self._generate_sticker_png(size=self._rand.randint(120, 180))
                except Exception:
                    continue
                pw = int(sw * self._rand.uniform(0.10, 0.18))
                ph = int(sh * self._rand.uniform(0.10, 0.18))
                left = self._rand.randint(0, max(0, sw - pw))
                top = self._rand.randint(0, max(0, sh - ph))
                try:
                    pic = slide.shapes.add_picture(path, left, top, width=pw, height=ph)
                    pic.rotation = self._rand.randint(-20, 20)
                except Exception:
                    pass
        except Exception as e:
            logger.debug(f"Could not insert stickers: {e}")

    def _set_background_picture(self, slide) -> None:
        """Set slide background to a generated neon pattern image."""
        try:
            sw_px, sh_px = 1600, 900
            img = Image.new('RGB', (sw_px, sh_px), (255, 255, 255))
            d = ImageDraw.Draw(img)
            # Diagonal stripes
            colors = ['#ff00ff','#00ffff','#ffff00','#ff0000','#00ff00','#ff8800']
            step = 40
            for i in range(-sh_px, sw_px, step):
                d.line([(i, 0), (i+sh_px, sh_px)], fill=self._rand.choice(colors), width=24)
            # Save temp
            with NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                img.save(tmp.name, format='PNG')
                path = tmp.name
            try:
                slide.background.fill.user_picture(path)
            except Exception:
                # Fallback: add a full-size picture
                slide.shapes.add_picture(path, 0, 0, width=self.presentation.slide_width, height=self.presentation.slide_height)
        except Exception as e:
            logger.debug(f"Could not set background picture: {e}")

    def _insert_gaudy_chart(self, slide) -> None:
        """Insert an intentionally unattractive chart with neon colors."""
        try:
            chart_data = CategoryChartData()
            chart_data.categories = ["A","B","C","D"]
            chart_data.add_series("Performance", [self._rand.randint(10,90) for _ in range(4)])
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(6)
            height = Inches(3.5)
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
            )
            chart = chart_shape.chart
            # Neon colors per data point
            series = chart.series[0]
            for i, point in enumerate(series.points):
                try:
                    point.format.fill.solid()
                    rgb = self._rand.choice(NEON_COLORS)
                    point.format.fill.fore_color.rgb = RGBColor(*rgb)
                except Exception:
                    continue
            # Loud chart title
            try:
                chart.chart_title.has_text_frame = True
                chart.chart_title.text_frame.text = "Super Important Numbers!!!"
            except Exception:
                pass
        except Exception as e:
            logger.debug(f"Could not insert chart: {e}")


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    # Example usage
    # presentation = Presentation("sample.pptx")
    # generator = TacoGenerator(presentation, tacky_level=7)
    # generator.apply_tacky_design()
    # presentation.save("sample_TACKY.pptx")
