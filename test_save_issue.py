#!/usr/bin/env python
"""Test presentation.save() directly to diagnose the issue"""

import os
import tempfile
from pptx import Presentation
from pptx.util import Inches

print("Testing python-pptx file saving...")

# Create a presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'Test'
subtitle.text = 'Test'

# Test 1: Save to temp directory (like app.py does)
print("\nTest 1: Saving to DasaMaker/uploads...")
os.makedirs('uploads', exist_ok=True)
output_path = 'uploads/test_output.pptx'

try:
    presentation.save(output_path)
    print(f"✗ ERROR: 'presentation' is not defined!")
except NameError:
    print(f"✓ Caught NameError - using correct variable 'prs'")

try:
    prs.save(output_path)
    if os.path.exists(output_path):
        file_size = os.path.getsize(output_path)
        print(f"✓ File saved successfully: {output_path}")
        print(f"✓ File size: {file_size} bytes")
        
        # Try to open it
        verify = Presentation(output_path)
        print(f"✓ File opens successfully: {len(verify.slides)} slides")
    else:
        print(f"✗ File not created at {output_path}")
except Exception as e:
    print(f"✗ Error saving: {e}")
    import traceback
    traceback.print_exc()

# Test 2: Check variable naming in the actual app code
print("\nTest 2: Checking app.py code for variable naming issues...")
with open('web/app.py', 'r') as f:
    content = f.read()
    # Look for variable name issues
    if 'presentation.save' in content:
        print("✓ Found 'presentation.save()' in app.py")
    else:
        print("✗ Could not find 'presentation.save()' in app.py")

# Test 3: BytesIO buffer test
print("\nTest 3: Testing BytesIO buffer approach...")
from io import BytesIO

prs2 = Presentation()
prs2.slide_width = Inches(10)
prs2.slide_height = Inches(7.5)
slide = prs2.slides.add_slide(prs2.slide_layouts[0])

# Save to memory buffer
buffer = BytesIO()
prs2.save(buffer)
buffer.seek(0)
buffer_data = buffer.read()

print(f"✓ Saved to BytesIO buffer: {len(buffer_data)} bytes")

# Save from buffer to file
with open('uploads/test_from_buffer.pptx', 'wb') as f:
    f.write(buffer_data)

if os.path.exists('uploads/test_from_buffer.pptx'):
    print(f"✓ Saved from buffer to file: {os.path.getsize('uploads/test_from_buffer.pptx')} bytes")
    verify2 = Presentation('uploads/test_from_buffer.pptx')
    print(f"✓ File from buffer opens successfully")
else:
    print(f"✗ Failed to save from buffer")

print("\n--- Diagnostic tests complete ---")
