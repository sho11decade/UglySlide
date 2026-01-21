#!/usr/bin/env python
"""Test each tacky level to identify which one causes corruption"""

import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from src.taco_generator import TacoGenerator

print("=== Testing each tacky level for file corruption ===\n")

test_file = 'samples/test_sample.pptx'

# Test each level from 1 to 10
for level in range(1, 11):
    print(f"Testing Level {level}...", end=" ")
    
    try:
        # Load original
        prs = Presentation(test_file)
        
        # Apply tacky design
        generator = TacoGenerator(prs, tacky_level=level, seed=42)
        generator.apply_tacky_design()
        
        # Save to temp file
        output_path = f'uploads/test_level_{level}.pptx'
        prs.save(output_path)
        
        # Verify by re-opening
        verify = Presentation(output_path)
        slides = len(verify.slides)
        file_size = os.path.getsize(output_path)
        
        print(f"✓ OK ({file_size} bytes, {slides} slides)")
        
    except Exception as e:
        print(f"✗ CORRUPTION: {e}")
        import traceback
        traceback.print_exc()
        break

print("\n--- Test complete ---")
print("If a level fails, that's the culprit causing the file corruption.")
