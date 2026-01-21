#!/usr/bin/env python
"""
Comprehensive test to identify exactly when corruption happens
"""

import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from src.ppt_analyzer import PPTAnalyzer
from src.taco_generator import TacoGenerator
from src.content_transformer import ContentTransformer

print("=== Complete API flow test with both design and content transformation ===\n")

test_file = 'samples/test_sample.pptx'

for level in range(1, 11):
    print(f"\nLevel {level}:")
    
    try:
        # Load original
        prs = Presentation(test_file)
        print(f"  1. Loaded: {len(prs.slides)} slides")
        
        # Apply design
        design_gen = TacoGenerator(prs, tacky_level=level, seed=42)
        design_gen.apply_tacky_design()
        print(f"  2. Design applied")
        
        # Apply content
        content_trans = ContentTransformer(prs, intensity=level, seed=42)
        content_trans.transform_all_content()
        print(f"  3. Content transformed")
        
        # Save
        output = f'uploads/full_level_{level}.pptx'
        prs.save(output)
        print(f"  4. Saved: {os.path.getsize(output)} bytes")
        
        # Verify
        verify = Presentation(output)
        print(f"  5. Verified: {len(verify.slides)} slides, ✓ OK")
        
        # Clean up objects
        del prs
        del design_gen
        del content_trans
        del verify
        
    except Exception as e:
        print(f"  ✗ FAILED at level {level}: {e}")
        import traceback
        traceback.print_exc()
        break

print("\n--- Test Complete ---")
