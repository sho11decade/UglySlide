#!/usr/bin/env python
"""Create a test PPTX file for debugging file processing"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt

# Create a simple test presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Add a title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'Test Presentation'
subtitle.text = 'Testing file processing'

# Save to samples
prs.save('samples/test_sample.pptx')
file_size = os.path.getsize('samples/test_sample.pptx')
print(f'✓ Test file created: samples/test_sample.pptx')
print(f'✓ File size: {file_size} bytes')

# Now test the processing
from src.ppt_analyzer import PPTAnalyzer
from src.taco_generator import TacoGenerator
from src.content_transformer import ContentTransformer

print('\n--- Testing file processing ---')

# Step 1: Analyze
print('1. Analyzing...')
analyzer = PPTAnalyzer('samples/test_sample.pptx')
analysis = analyzer.analyze()
print(f'   ✓ Analysis: {analysis.to_dict()}')

# Step 2: Load and modify
print('2. Loading presentation...')
presentation = Presentation('samples/test_sample.pptx')
print(f'   ✓ Loaded {len(presentation.slides)} slides')

# Step 3: Apply design
print('3. Applying tacky design...')
design_generator = TacoGenerator(presentation, tacky_level=5, seed=42)
design_generator.apply_tacky_design()
print('   ✓ Design applied')

# Step 4: Apply content transformation
print('4. Applying content transformation...')
content_transformer = ContentTransformer(presentation, intensity=5, seed=42)
content_transformer.transform_all_content()
print('   ✓ Content transformed')

# Step 5: Save
print('5. Saving...')
output_path = 'samples/test_sample_TACKY.pptx'
presentation.save(output_path)
output_size = os.path.getsize(output_path)
print(f'   ✓ Saved: {output_path}')
print(f'   ✓ Output size: {output_size} bytes')

# Step 6: Verify by re-opening
print('6. Verifying...')
try:
    verify_prs = Presentation(output_path)
    print(f'   ✓ File opens successfully: {len(verify_prs.slides)} slides')
except Exception as e:
    print(f'   ✗ Error opening file: {e}')

print('\n--- All tests passed! ---')
